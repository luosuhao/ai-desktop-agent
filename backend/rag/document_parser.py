"""Multi-format Document Parser

Supports: PDF, Word, PPT, Excel, CSV, Markdown, TXT, Images
"""

import os
import re
import json
import csv
from typing import List, Dict, Optional, Any
from datetime import datetime
from io import BytesIO


class ParsedDocument:
    """Represents a parsed document with structured content"""
    def __init__(self, doc_id: str, filename: str):
        self.id = doc_id
        self.filename = filename
        self.file_type = os.path.splitext(filename)[1].lower()
        self.pages: List[Dict] = []  # {page_num, text, tables, images, formulas}
        self.tables: List[Dict] = []  # {page_num, rows, cols, headers, data, markdown}
        self.images: List[Dict] = []  # {page_num, caption, path}
        self.metadata: Dict = {
            "title": "",
            "author": "",
            "created": "",
            "page_count": 0,
            "parse_time": ""
        }
        self.chunks: List[Dict] = []  # For RAG indexing
        self.parse_errors: List[str] = []


class DocumentParser:
    """Parse various document formats into structured content"""

    @staticmethod
    def parse(file_path: str, doc_id: str) -> ParsedDocument:
        """Parse document based on file extension"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        doc = ParsedDocument(doc_id, filename)

        start_time = datetime.now()

        try:
            if ext == '.pdf':
                DocumentParser._parse_pdf(file_path, doc)
            elif ext == '.docx':
                DocumentParser._parse_docx(file_path, doc)
            elif ext == '.pptx':
                DocumentParser._parse_pptx(file_path, doc)
            elif ext in ('.xlsx', '.xls'):
                DocumentParser._parse_excel(file_path, doc)
            elif ext == '.csv':
                DocumentParser._parse_csv(file_path, doc)
            elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp'):
                DocumentParser._parse_image(file_path, doc)
            elif ext in ('.md', '.txt'):
                DocumentParser._parse_text(file_path, doc)
            else:
                doc.parse_errors.append(f"Unsupported file type: {ext}")
        except Exception as e:
            doc.parse_errors.append(f"Parse error: {str(e)}")

        doc.metadata["parse_time"] = str(datetime.now() - start_time)
        doc.metadata["page_count"] = len(doc.pages)

        # Generate chunks for RAG
        DocumentParser._generate_chunks(doc)

        return doc

    @staticmethod
    def _parse_pdf(file_path: str, doc: ParsedDocument):
        """Parse PDF file using pdfminer"""
        try:
            from pdfminer.high_level import extract_pages
            from pdfminer.layout import LTTextBox, LTTable, LTFigure, LAParams, LTChar

            laparams = LAParams(detect_vertical=True, all_texts=True)
            pages = list(extract_pages(file_path, laparams=laparams))

            for page_num, page_layout in enumerate(pages, 1):
                page_text = ""
                page_tables = []

                for element in page_layout:
                    if isinstance(element, LTTextBox):
                        page_text += element.get_text() + "\n"

                doc.pages.append({
                    "page_num": page_num,
                    "text": page_text.strip(),
                    "tables_count": len(page_tables),
                    "images_count": 0
                })

                # Simple table detection by looking for tabular patterns
                lines = page_text.split('\n')
                table_lines = [l for l in lines if l.count('\t') > 1 or l.count('  ') > 2]
                if table_lines:
                    doc.tables.append({
                        "page_num": page_num,
                        "source": "pdf_text",
                        "rows": len(table_lines),
                        "cols": max(l.count('\t') + 1 if '\t' in l else 3 for l in table_lines),
                        "data": [re.split(r'\t|\s{2,}', l.strip()) for l in table_lines],
                        "markdown": DocumentParser._to_markdown(table_lines)
                    })

        except ImportError:
            # Fallback to PyPDF2
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text() or ""
                    doc.pages.append({
                        "page_num": page_num,
                        "text": text,
                        "tables_count": 0,
                        "images_count": 0
                    })
                doc.metadata["page_count"] = len(reader.pages)
            except ImportError:
                doc.parse_errors.append("No PDF parser available (install pdfminer.six or PyPDF2)")

    @staticmethod
    def _parse_docx(file_path: str, doc: ParsedDocument):
        """Parse Word document"""
        from docx import Document as DocxDocument
        docx = DocxDocument(file_path)

        current_text = []
        for para in docx.paragraphs:
            current_text.append(para.text)

        # Extract tables
        for table in docx.tables:
            rows_data = []
            headers = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if i == 0:
                    headers = cells
                rows_data.append(cells)

            doc.tables.append({
                "page_num": 1,
                "source": "docx",
                "rows": len(rows_data),
                "cols": len(headers) if headers else max(len(r) for r in rows_data) if rows_data else 0,
                "headers": headers,
                "data": rows_data,
                "markdown": DocumentParser._to_markdown_table(headers, rows_data[1:] if headers else rows_data)
            })

        doc.pages.append({
            "page_num": 1,
            "text": '\n'.join(current_text),
            "tables_count": len(doc.tables),
            "images_count": 0
        })

        # Metadata
        if docx.core_properties:
            doc.metadata["title"] = docx.core_properties.title or ""
            doc.metadata["author"] = docx.core_properties.author or ""

    @staticmethod
    def _parse_pptx(file_path: str, doc: ParsedDocument):
        """Parse PowerPoint document"""
        from pptx import Presentation
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    rows_data = []
                    headers = []
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        if i == 0:
                            headers = cells
                        rows_data.append(cells)
                    doc.tables.append({
                        "page_num": slide_num,
                        "source": "pptx",
                        "rows": len(rows_data),
                        "cols": len(headers) if headers else max(len(r) for r in rows_data) if rows_data else 0,
                        "headers": headers,
                        "data": rows_data,
                        "markdown": DocumentParser._to_markdown_table(headers, rows_data[1:] if headers else rows_data)
                    })

            doc.pages.append({
                "page_num": slide_num,
                "text": '\n'.join(slide_text),
                "tables_count": len(doc.tables),
                "images_count": 0
            })

    @staticmethod
    def _parse_excel(file_path: str, doc: ParsedDocument):
        """Parse Excel file"""
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data = []
            headers = []

            for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if row_idx == 1:
                    headers = cells
                rows_data.append({'row': row_idx, 'cells': cells})

            doc.tables.append({
                "page_num": 1,
                "source": f"excel:{sheet_name}",
                "sheet_name": sheet_name,
                "rows": len(rows_data),
                "cols": len(headers),
                "headers": headers,
                "data": [r['cells'] for r in rows_data],
                "markdown": DocumentParser._to_markdown_table(headers, rows_data[1:])
            })

            # Page content for the sheet
            all_text = '\n'.join([f"Sheet: {sheet_name}"] +
                                 ['\t'.join(r['cells']) for r in rows_data])
            doc.pages.append({
                "page_num": 1,
                "text": all_text,
                "tables_count": 1,
                "images_count": 0,
                "sheet_name": sheet_name
            })

        doc.metadata["page_count"] = len(wb.sheetnames)

    @staticmethod
    def _parse_csv(file_path: str, doc: ParsedDocument):
        """Parse CSV file"""
        encoding = 'utf-8'
        # Try to detect encoding
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                first_bytes = f.read(1000)
            # Check for BOM
            if first_bytes.startswith('﻿'):
                encoding = 'utf-8-sig'
        except UnicodeDecodeError:
            encoding = 'gbk'

        with open(file_path, 'r', encoding=encoding) as f:
            reader = csv.reader(f)
            rows_data = []
            headers = []
            for i, row in enumerate(reader):
                if i == 0:
                    headers = row
                rows_data.append(row)

        doc.tables.append({
            "page_num": 1,
            "source": "csv",
            "rows": len(rows_data),
            "cols": len(headers),
            "headers": headers,
            "data": rows_data,
            "markdown": DocumentParser._to_markdown_table(headers, rows_data[1:])
        })

        all_text = '\n'.join([','.join(r) for r in rows_data])
        doc.pages.append({
            "page_num": 1,
            "text": all_text,
            "tables_count": 1,
            "images_count": 0
        })

    @staticmethod
    def _parse_image(file_path: str, doc: ParsedDocument):
        """Parse image file with OCR"""
        try:
            from PIL import Image
            import pytesseract

            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')

            # Try to detect tables in the image
            try:
                import pytesseract
                table_data = pytesseract.image_to_data(img, lang='chi_sim+eng', output_type=pytesseract.Output.DICT)
                # Simple table detection by checking for aligned text
                lines = {}
                for i, text_content in enumerate(table_data['text']):
                    if text_content.strip():
                        line_num = table_data['line_num'][i]
                        if line_num not in lines:
                            lines[line_num] = []
                        lines[line_num].append(text_content.strip())

                table_lines = [lines[k] for k in sorted(lines.keys())]
                if len(table_lines) > 1 and all(len(tl) > 1 for tl in table_lines):
                    max_cols = max(len(tl) for tl in table_lines)
                    if max_cols >= 2:
                        doc.tables.append({
                            "page_num": 1,
                            "source": "ocr_image",
                            "rows": len(table_lines),
                            "cols": max_cols,
                            "data": table_lines,
                            "markdown": DocumentParser._to_markdown_table(
                                table_lines[0], table_lines[1:]
                            )
                        })
            except Exception as e:
                doc.parse_errors.append(f"Image table detection note: {str(e)}")

            doc.pages.append({
                "page_num": 1,
                "text": text,
                "tables_count": len(doc.tables),
                "images_count": 1
            })

        except ImportError as e:
            doc.parse_errors.append(f"Image/OCR library missing: {str(e)}")
            doc.pages.append({"page_num": 1, "text": "[Image file - OCR not available]", "tables_count": 0, "images_count": 1})

    @staticmethod
    def _parse_text(file_path: str, doc: ParsedDocument):
        """Parse plain text or markdown file"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        content = None

        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    content = f.read()
                break
            except (UnicodeDecodeError, UnicodeError):
                continue

        if content is None:
            content = "[Unable to read file content]"

        # Simple Markdown table detection
        table_lines = []
        in_table = False
        for line in content.split('\n'):
            if '|' in line and line.count('|') >= 2:
                table_lines.append(line)
                in_table = True
            elif in_table and not line.strip():
                in_table = False

        if table_lines:
            headers = [h.strip() for h in table_lines[0].split('|') if h.strip()]
            data_rows = []
            for line in table_lines[2:]:  # Skip separator line
                cells = [c.strip() for c in line.split('|') if c.strip()]
                if cells:
                    data_rows.append(cells)
            doc.tables.append({
                "page_num": 1,
                "source": "markdown",
                "rows": len(data_rows),
                "cols": len(headers),
                "headers": headers,
                "data": data_rows,
                "markdown": DocumentParser._to_markdown_table(headers, data_rows)
            })

        doc.pages.append({
            "page_num": 1,
            "text": content,
            "tables_count": len(doc.tables),
            "images_count": 0
        })

    @staticmethod
    def _to_markdown_table(headers: List[str], data_rows: List[List[str]]) -> str:
        """Convert headers and rows to Markdown table"""
        if not headers and not data_rows:
            return ""

        if not headers and data_rows:
            headers = data_rows[0]
            data_rows = data_rows[1:]

        md = "| " + " | ".join(headers) + " |\n"
        md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        for row in data_rows:
            md += "| " + " | ".join(row[:len(headers)]) + " |\n"
        return md

    @staticmethod
    def _to_markdown(lines: List[str]) -> str:
        """Convert lines to a simple markdown representation"""
        def split_line(l: str) -> str:
            parts = re.split(r'\t|\s{2,}', l.strip())
            return ' | '.join(parts)
        return '\n'.join([f"| {split_line(l)} |" for l in lines])

    @staticmethod
    def _generate_chunks(doc: ParsedDocument, chunk_size: int = 500, overlap: int = 50):
        """Generate text chunks for RAG indexing"""
        doc.chunks = []

        # Chunk from pages
        for page in doc.pages:
            text = page.get("text", "")
            if not text:
                continue

            words = text.split()
            for i in range(0, len(words), chunk_size - overlap):
                chunk_text = ' '.join(words[i:i + chunk_size])
                if chunk_text.strip():
                    doc.chunks.append({
                        "id": f"{doc.id}_p{page['page_num']}_c{len(doc.chunks)}",
                        "content": chunk_text,
                        "chunk_type": "text",
                        "page_number": page['page_num'],
                        "metadata": {
                            "document_id": doc.id,
                            "filename": doc.filename,
                            "file_type": doc.file_type
                        }
                    })

        # Chunk from tables
        for table in doc.tables:
            md = table.get("markdown", "")
            if md:
                doc.chunks.append({
                    "id": f"{doc.id}_t{len(doc.chunks)}",
                    "content": md,
                    "chunk_type": "table",
                    "page_number": table.get("page_num", 1),
                    "metadata": {
                        "document_id": doc.id,
                        "filename": doc.filename,
                        "file_type": doc.file_type,
                        "table_source": table.get("source", ""),
                        "rows": table.get("rows", 0),
                        "cols": table.get("cols", 0)
                    }
                })
